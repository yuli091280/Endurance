from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as MSO_SHAPE
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT as PP_ALIGN
from pptx.enum.dml import MSO_THEME_COLOR_INDEX as MSO_THEME_COLOR

BUTTON_ROW_HEIGHT = Inches(6.75)
BUTTON_SPACING = Inches(0.2)
BUTTON_WIDTH = Inches(0.55)
BUTTON_HEIGHT = Inches(0.5)

TABLE_ROW_HEIGHT = Inches(2)
TABLE_LEFT_PADDING = Inches(0.5)
TABLE_WIDTH = Inches(9)
TABLE_HEIGHT = Inches(0.5)
TABLE_FONT_SIZE = Pt(10)

MIDDLE_OF_SLIDE = Presentation().slide_width / 2

MAX_ROWS_PER_SLIDE = 15
MAX_BUTTONS_PER_SLIDE = 10
MAX_BTNS_PER_SLIDE_WITH_ARROWS = MAX_BUTTONS_PER_SLIDE + 2


def _get_proper_button_spacing(button_num, total_buttons):
    """
    Places a button according to the total number of buttons on the slide.

    :param button_num: Number of button being placed
    :type button_num: int
    :param total_buttons: Total number of buttons on the slide
    :type total_buttons: int
    :return: Location to place button in Inches
    :rtype: pptx.util.Inches
    """
    if total_buttons % 2 == 0:
        # Even buttons, no center button
        middle_button = total_buttons // 2
        if button_num < middle_button:
            # Button is on left side of the slide
            return (
                MIDDLE_OF_SLIDE
                + (BUTTON_SPACING / 2)
                + (button_num - middle_button - 1) * (BUTTON_WIDTH + BUTTON_SPACING)
            )
        else:
            # Button is on the right side of the slide
            return (
                MIDDLE_OF_SLIDE
                - (BUTTON_SPACING / 2)
                - BUTTON_WIDTH
                - (middle_button - button_num) * (BUTTON_WIDTH + BUTTON_SPACING)
            )
    else:
        # Odd buttons, there is a middle button
        middle_button = (total_buttons // 2) + 1
        if button_num < middle_button:  # Button is on left side of the slide
            return (
                MIDDLE_OF_SLIDE
                + (BUTTON_WIDTH / 2)
                + BUTTON_SPACING
                + (button_num - middle_button - 1) * (BUTTON_WIDTH + BUTTON_SPACING)
            )
        else:
            # Button is on right side of the slide
            return (
                MIDDLE_OF_SLIDE
                - (BUTTON_WIDTH / 2)
                - (middle_button - button_num) * (BUTTON_WIDTH + BUTTON_SPACING)
            )


def _button_factory(total_buttons, slide, button_place, text, color=None):
    """
    Generate a button shape using the specified parameters.

    :param total_buttons: Total number of buttons on the slide
    :type total_buttons: int
    :param slide: Slide to place button on
    :type slide: pptx.slide.Slide
    :param button_place: Location on slide to put button
    :type button_place: int
    :param text: Text content of button
    :type text: str
    :param color: Optional color of button
    :type color: pptx.enum.dml.MSO_THEME_COLOR_INDEX
    :return: Formatted button
    :rtype: pptx.shapes.autoshapes.Shape
    """
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        _get_proper_button_spacing(button_place, total_buttons),
        BUTTON_ROW_HEIGHT,
        BUTTON_WIDTH,
        BUTTON_HEIGHT,
    )
    shape.text_frame.text = text
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    if color:
        shape.fill.gradient()
        for gradient_stop in shape.fill.gradient_stops:
            gradient_stop.color.theme_color = color
        shape.line.color.theme_color = color

    return shape


def generate_powerpoint(selected_query, data, headers, file_path):
    """
    Generate a PowerPoint using the supplied header and table data, saving to the specified file path.

    :param selected_query: The title of the query being displayed
    :type selected_query: str
    :param data: The data of the query being displayed
    :type data: list[tuple]
    :param headers: The headers of the data being displayed
    :type headers: list[str]
    :param file_path: The file path to save the PowerPoint to
    :type file_path: str
    """
    page_num = 1
    all_buttons = []

    # Create a PowerPoint presentation
    presentation = Presentation()

    # Determine the number of chunks needed based on the maximum rows per slide
    num_rows = len(data)
    chunks = [
        data[i : i + MAX_ROWS_PER_SLIDE] for i in range(0, num_rows, MAX_ROWS_PER_SLIDE)
    ]

    num_slides = len(chunks)

    # Number of slides plus forward/backward
    num_buttons = num_slides + 2

    slide_layout = presentation.slide_layouts[5]
    for index, chunk in enumerate(chunks):
        slide_buttons = []

        # Use the layout that supports title and content
        slide = presentation.slides.add_slide(slide_layout)
        current_slide = index + 1

        # Set table title
        if num_slides > 1:
            slide.shapes.title.text = f"{selected_query} Page {page_num}"
        else:
            slide.shapes.title.text = selected_query

        # Add page number to slides
        page_number_text = slide.shapes.add_textbox(0, 0, Inches(1), Inches(1))
        page_number_text.text_frame.text = str(page_num)
        page_num += 1

        # Add a table to the slide
        table = slide.shapes.add_table(
            rows=len(chunk) + 1,
            cols=len(headers),
            left=TABLE_LEFT_PADDING,
            top=TABLE_ROW_HEIGHT,
            width=TABLE_WIDTH,
            height=TABLE_HEIGHT,
        ).table

        # Set column names as the first row of the table
        for col_num, col_name in enumerate(headers):
            cell = table.cell(0, col_num)
            cell.text = col_name
            cell.text_frame.paragraphs[0].font.size = TABLE_FONT_SIZE

        # Populate the table with data
        for row_num in range(len(chunk)):
            for col_num, value in enumerate(chunk[row_num]):
                cell = table.cell(row_num + 1, col_num)
                cell.text = str(value)
                cell.text_frame.paragraphs[0].font.size = TABLE_FONT_SIZE

        # If there's only one slide, skip buttons
        if num_slides == 1:
            continue

        # Create factory function with pre-filled slide and button info
        create_button = lambda button_num, text, color=None: _button_factory(
            min(num_buttons, MAX_BTNS_PER_SLIDE_WITH_ARROWS),
            slide,
            button_num,
            text,
            color,
        )

        # Add previous slide button to slide
        slide_buttons.append(create_button(1, "<", MSO_THEME_COLOR.ACCENT_6))

        # Add numbered slide buttons to slide
        if num_buttons > MAX_BTNS_PER_SLIDE_WITH_ARROWS:
            # Not all the buttons will fit, truncate
            if current_slide <= 4:
                # We're at the start of the list, include the first N slides
                button_range = range(1, MAX_BUTTONS_PER_SLIDE + 1)
            elif current_slide >= num_slides - 5:
                # We're at the end of the list, include the last N slides
                button_range = range(
                    num_slides - MAX_BUTTONS_PER_SLIDE + 1, num_slides + 1
                )
            else:
                # We're somewhere in the middle, include the 4 slides before and the 5 slides after this one
                button_range = range(current_slide - 4, current_slide + 6)

            # Button position is decoupled from the number, skipping previous slide button
            button_pos = 2
            for button in button_range:
                slide_buttons.append(
                    create_button(
                        button_pos,
                        str(button),
                        # Add special formatting for selected slide
                        MSO_THEME_COLOR.ACCENT_2 if button == current_slide else None,
                    )
                )
                button_pos += 1
        else:
            # All the buttons will fit, we can use their number to decide their position
            for button in range(1, num_slides + 1):
                slide_buttons.append(
                    create_button(
                        button + 1,
                        str(button),
                        # Add special formatting for selected slide
                        MSO_THEME_COLOR.ACCENT_2 if button == current_slide else None,
                    )
                )

        # Add next slide button to slide
        slide_buttons.append(
            create_button(
                min(num_buttons, MAX_BTNS_PER_SLIDE_WITH_ARROWS),
                ">",
                MSO_THEME_COLOR.ACCENT_6,
            )
        )

        all_buttons.append(slide_buttons)

    # After our slides have been created, add functionality to buttons
    for slide_index, slide_button_list in enumerate(all_buttons):
        # Add functionality to previous slide button
        slide_button_list[0].click_action.target_slide = presentation.slides[
            slide_index - 1
        ]

        # Add functionality to next slide button
        slide_button_list[-1].click_action.target_slide = presentation.slides[
            (slide_index + 1) % num_slides
        ]

        # Add functionality to rest of buttons
        for button in slide_button_list[1:-1]:
            button.click_action.target_slide = presentation.slides[
                int(button.text_frame.text) - 1
            ]

    # Save the PowerPoint presentation
    presentation.save(file_path)
